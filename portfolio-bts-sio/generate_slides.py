
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Helper to set background quickly
    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    # Initialize
    prs = Presentation()
    
    # Theme Colors
    DARK_BG = RGBColor(20, 20, 30)       # Dark Blue-Grey
    ACCENT_ORANGE = RGBColor(255, 100, 50) # Vibrant Orange
    TEXT_WHITE = RGBColor(240, 240, 240)
    TEXT_GREY = RGBColor(180, 180, 190)

    # ---------------------------------------------------------
    # Helper Functions
    # ---------------------------------------------------------
    
    def add_slide_title_only(title_text, subtitle_text=""):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, DARK_BG)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        # Style Title
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_ORANGE
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(44)
        
        subtitle.text = subtitle_text
        # Style Subtitle
        if subtitle.text_frame.paragraphs:
            for p in subtitle.text_frame.paragraphs:
                p.font.color.rgb = TEXT_WHITE
                p.font.size = Pt(24)

    def add_slide_content(title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, DARK_BG)
        
        title = slide.shapes.title
        title.text = title_text
        # Style Title
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_ORANGE
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        
        # Content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear() 
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.color.rgb = TEXT_WHITE
            p.font.size = Pt(20)
            p.space_after = Pt(14)
            
            # Check for sub-items (simple indent logic for string starting with space)
            if item.startswith("   -"):
                p.level = 1

    # ---------------------------------------------------------
    # Slides Generation
    # ---------------------------------------------------------

    # Slide 1: Title
    add_slide_title_only("Épreuve E5 - Support et Mise à Disposition", 
                         "BTS SIO - Option SLAM\nCandidat : Benoit DUFOUR - Session 2026")

    # Slide 2: Sommaire
    add_slide_content("Sommaire", [
        "1. Présentation du candidat",
        "2. Contexte professionnel (Stages)",
        "3. Gestion du patrimoine informatique",
        "4. Projet Phare : Gestion de Tickets",
        "5. Projet Secondaire : Déploiement & Maintenance",
        "6. Conclusion"
    ])

    # Slide 3: Présentation
    add_slide_content("Présentation du Candidat", [
        "Benoit DUFOUR",
        "Formation : BTS SIO Option SLAM",
        "   - Solutions Logicielles et Applications Métiers",
        "Parcours :",
        "   - Bac Général (NSI / Maths)",
        "   - Autodidacte passionné",
        "Objectifs :",
        "   - Devenir Développeur Full Stack",
        "   - Poursuite d'études (Licence / École d'Ingé)"
    ])

    # Slide 4: Contexte Pro
    add_slide_content("Contexte Professionnel", [
        "Stage 2ème Année - Entreprise ABC (6 semaines)",
        "   - Développement & Maintenance App Web Interne",
        "   - Environnement : Symfony, Git, Scrum",
        "",
        "Stage 1ère Année - Entreprise XYZ (5 semaines)",
        "   - Création Site Vitrine & Support Utilisateur",
        "   - Environnement : CMS Wordpress, HTML/CSS"
    ])

    # Slide 5: Patrimoine
    add_slide_content("Gestion du Patrimoine Informatique", [
        "Outils de Versionning :",
        "   - Usage quotidien de Git & GitHub",
        "Gestion des Incidents :",
        "   - Flux de tickets (GLPI / Jira)",
        "Documentation :",
        "   - Rédaction de docs techniques et guides utilisateurs",
        "Sécurité & Veille :",
        "   - Certifications ANSSI & Cisco Cybersecurity"
    ])

    # Slide 6: Projet 1
    add_slide_content("Projet Phare : Gestion de Tickets", [
        "Contexte :",
        "   - Besoin d'une solution centralisée pour les incidents",
        "Service Rendu (E5) :",
        "   - Interface de création et suivi de tickets",
        "   - Tableau de bord technicien",
        "Mise à disposition :",
        "   - Application Web (PHP/MySQL)",
        "   - Déploiement serveur local (WAMP/XAMPP)"
    ])

    # Slide 7: Projet 2
    add_slide_content("Projet Secondaire : Déploiement", [
        "Projet : Site E-Commerce (Laravel)",
        "Actions de Mise à Disposition :",
        "   - Configuration de l'hébergement Web (Apache)",
        "   - Import et sécurisation de la Base de Données",
        "   - Gestion des environnements (.env)",
        "Maintenance :",
        "   - Sauvegardes et mises à jour de sécurité"
    ])

    # Slide 8: Conclusion
    add_slide_content("Conclusion", [
        "Bilan :",
        "   - Maîtrise du cycle de vie des services",
        "   - Compétences techniques validées",
        "Perspectives :",
        "   - Approfondissement DevOps (Docker, CI/CD)",
        "   - Veille continue sur l'IA"
    ])

    # Slide 9: Fin
    add_slide_title_only("Merci de votre attention", "Avez-vous des questions ?")

    # Save
    filename = 'presentation_e5_benoit_dufour.pptx'
    prs.save(filename)
    print(f"Presentation saved to {filename}")

if __name__ == "__main__":
    create_presentation()
